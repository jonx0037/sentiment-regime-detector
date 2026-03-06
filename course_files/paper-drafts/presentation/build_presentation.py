"""
Build SMU Capstone Presentation using the dark SMU template.
Uses python-pptx to create slides from template layouts, fill content, and add images.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ── Paths ──
BASE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = "/Users/jonathanrocha/Documents/SMU/DS_6210_Capstone/course_files/course_artifacts/smu_capstone_template_-_dark.pptx"
IMG_DIR = BASE  # figures and screenshots are in presentation/
OUTPUT = os.path.join(BASE, "JonathanRocha_Capstone_Presentation.pptx")

# ── SMU Colors ──
SMU_BLUE = RGBColor(0x02, 0x57, 0xA1)
SMU_RED = RGBColor(0xC0, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
GOLD = RGBColor(0xFF, 0xC0, 0x00)

# ── Load template ──
prs = Presentation(TEMPLATE)

# Map layouts by name for clarity
layout_map = {}
for layout in prs.slide_layouts:
    layout_map[layout.name] = layout
print("Available layouts:", list(layout_map.keys()))

# Delete the instruction slide (slide 0)
rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
prs.part.drop_rel(rId)
del prs.slides._sldIdLst[0]

# ── Helper functions ──
def add_slide(layout_name):
    """Add a slide with the given layout name."""
    layout = layout_map[layout_name]
    return prs.slides.add_slide(layout)

def set_title(slide, text):
    """Set the title placeholder text."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:  # title
            ph.text = text
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = WHITE
            return ph
    return None

def set_body(slide, idx=1):
    """Get body placeholder by index."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None

def add_bullet(tf, text, level=0, bold=False, font_size=14, color=WHITE):
    """Add a bullet point to a text frame."""
    p = tf.add_paragraph()
    p.level = level
    p.space_after = Pt(4)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p

def clear_body(ph):
    """Clear all paragraphs from placeholder and return text_frame."""
    tf = ph.text_frame
    tf.clear()
    return tf

def remove_placeholder(slide, idx):
    """Remove a placeholder by index (hide empty placeholders)."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            sp = ph._element
            sp.getparent().remove(sp)
            return True
    return False

def add_image_to_slide(slide, img_path, left, top, width=None, height=None):
    """Add an image to a slide at specified position."""
    full_path = os.path.join(IMG_DIR, img_path)
    if not os.path.exists(full_path):
        print(f"  WARNING: Image not found: {full_path}")
        return None
    if width and height:
        return slide.shapes.add_picture(full_path, left, top, width, height)
    elif width:
        return slide.shapes.add_picture(full_path, left, top, width=width)
    elif height:
        return slide.shapes.add_picture(full_path, left, top, height=height)
    return slide.shapes.add_picture(full_path, left, top)

def add_textbox(slide, left, top, width, height, text, font_size=12, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox

# ══════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════════════════
print("Creating Slide 1: Title")
s1 = add_slide("White")
for ph in s1.placeholders:
    if ph.placeholder_format.idx == 0:  # center title
        ph.text = "Cross-Asset Sentiment\nRegime Detector"
        for para in ph.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size = Pt(36)
                run.font.bold = True
                run.font.color.rgb = SMU_BLUE
    elif ph.placeholder_format.idx == 1:  # subtitle
        tf = ph.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = "Automating Market Psychology Analysis Through Multi-Source NLP"
        r1.font.size = Pt(18)
        r1.font.color.rgb = RGBColor(0x44, 0x54, 0x6A)

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(20)
        r2 = p2.add_run()
        r2.text = "Jonathan Rocha"
        r2.font.size = Pt(16)
        r2.font.bold = True
        r2.font.color.rgb = SMU_RED

        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = "SMU Master of Science in Data Science  |  Capstone Project  |  Spring 2026"
        r3.font.size = Pt(12)
        r3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

# ══════════════════════════════════════════════════════════════
# SLIDE 2: Problem & Motivation
# ══════════════════════════════════════════════════════════════
print("Creating Slide 2: Problem & Motivation")
s2 = add_slide("Content")
set_title(s2, "Problem & Motivation")
body2 = set_body(s2, 1)
tf2 = clear_body(body2)

p = tf2.paragraphs[0]
r = p.add_run()
r.text = "The Challenge"
r.font.size = Pt(18)
r.font.bold = True
r.font.color.rgb = SMU_BLUE

add_bullet(tf2, "Financial markets generate massive unstructured text daily — news, social media, analyst reports", font_size=14)
add_bullet(tf2, "Sentiment shifts across asset classes often precede regime transitions", font_size=14)
add_bullet(tf2, "Manual monitoring is impossible at scale across multiple markets simultaneously", font_size=14)

p_gap = tf2.add_paragraph()
p_gap.space_before = Pt(12)
r_gap = p_gap.add_run()
r_gap.text = "Research Gap"
r_gap.font.size = Pt(18)
r_gap.font.bold = True
r_gap.font.color.rgb = SMU_BLUE

add_bullet(tf2, "Existing approaches focus on single-asset sentiment — no cross-asset regime detection", font_size=14)
add_bullet(tf2, "No unified framework integrating NLP sentiment with volatility regime models", font_size=14)

p_goal = tf2.add_paragraph()
p_goal.space_before = Pt(12)
r_goal = p_goal.add_run()
r_goal.text = "Goal: "
r_goal.font.size = Pt(14)
r_goal.font.bold = True
r_goal.font.color.rgb = GOLD
r_goal2 = p_goal.add_run()
r_goal2.text = "Build an automated system that detects market regime transitions using cross-asset NLP sentiment signals"
r_goal2.font.size = Pt(14)
r_goal2.font.color.rgb = WHITE

# ══════════════════════════════════════════════════════════════
# SLIDE 3: Data & NLP Pipeline (with dashboard sentiment screenshot)
# ══════════════════════════════════════════════════════════════
print("Creating Slide 3: Data & NLP Pipeline")
s3 = add_slide("OBJECT_WITH_CAPTION_TEXT")
set_title(s3, "Data & NLP Pipeline")

body3L = set_body(s3, 2)
tf3 = clear_body(body3L)

p = tf3.paragraphs[0]
r = p.add_run()
r.text = "~33 Million Documents"
r.font.size = Pt(14)
r.font.bold = True
r.font.color.rgb = SMU_BLUE

add_bullet(tf3, "Equities (SPY) — News & earnings", font_size=11)
add_bullet(tf3, "Crypto (BTC) — Social media & Reddit", font_size=11)
add_bullet(tf3, "Forex (EUR/USD) — Central bank comms", font_size=11)
add_bullet(tf3, "Commodities (Gold) — Reports & news", font_size=11)

p_pipe = tf3.add_paragraph()
p_pipe.space_before = Pt(10)
r_pipe = p_pipe.add_run()
r_pipe.text = "NLP Pipeline"
r_pipe.font.size = Pt(14)
r_pipe.font.bold = True
r_pipe.font.color.rgb = SMU_BLUE

add_bullet(tf3, "FinBERT ensemble transformer", font_size=11)
add_bullet(tf3, "Daily aggregation per asset", font_size=11)
add_bullet(tf3, "4,490 observations (2005–2025)", font_size=11)

p_feat = tf3.add_paragraph()
p_feat.space_before = Pt(8)
r_feat = p_feat.add_run()
r_feat.text = "→  22 engineered features"
r_feat.font.size = Pt(11)
r_feat.font.color.rgb = GOLD

remove_placeholder(s3, 1)
add_image_to_slide(s3, "dashboard_sentiment.png",
                   Inches(5.67), Inches(1.08), width=Inches(6.75))

# ══════════════════════════════════════════════════════════════
# SLIDE 4: Two-Layer Architecture (with dashboard GARCH/regime screenshot)
# ══════════════════════════════════════════════════════════════
print("Creating Slide 4: Two-Layer Architecture")
s4 = add_slide("OBJECT_WITH_CAPTION_TEXT")
set_title(s4, "Two-Layer Architecture")

body4L = set_body(s4, 2)
tf4L = clear_body(body4L)

p = tf4L.paragraphs[0]
r = p.add_run()
r.text = "Layer 1: GARCH(1,1)"
r.font.size = Pt(13)
r.font.bold = True
r.font.color.rgb = SMU_BLUE

add_bullet(tf4L, "Conditional variance per asset", font_size=11)
add_bullet(tf4L, "Volatility clustering + persistence", font_size=11)

p2 = tf4L.add_paragraph()
p2.space_before = Pt(8)
r2 = p2.add_run()
r2.text = "Layer 2: Jump Model"
r2.font.size = Pt(13)
r2.font.bold = True
r2.font.color.rgb = SMU_BLUE

add_bullet(tf4L, "Data-driven jump penalty (CV)", font_size=11)
add_bullet(tf4L, "Detects structural regime breaks", font_size=11)

p3 = tf4L.add_paragraph()
p3.space_before = Pt(8)
r3 = p3.add_run()
r3.text = "4 Regime States"
r3.font.size = Pt(13)
r3.font.bold = True
r3.font.color.rgb = GOLD

add_bullet(tf4L, "Low Volatility", font_size=10, color=RGBColor(0x2C, 0xA0, 0x2C))
add_bullet(tf4L, "Normal", font_size=10, color=RGBColor(0x1F, 0x77, 0xB4))
add_bullet(tf4L, "Elevated", font_size=10, color=RGBColor(0xFF, 0x7F, 0x0E))
add_bullet(tf4L, "High Volatility", font_size=10, color=SMU_RED)

p_net = tf4L.add_paragraph()
p_net.space_before = Pt(6)
r_net = p_net.add_run()
r_net.text = "Network: Transfer Entropy for cross-asset information flow (H3)"
r_net.font.size = Pt(9)
r_net.font.italic = True
r_net.font.color.rgb = LIGHT_GRAY

remove_placeholder(s4, 1)
add_image_to_slide(s4, "dashboard_garch_regime.png",
                   Inches(5.67), Inches(1.08), width=Inches(6.75))

# ══════════════════════════════════════════════════════════════
# SLIDE 5: Validation Framework (with COVID figure)
# ══════════════════════════════════════════════════════════════
print("Creating Slide 5: Validation Framework")
s5 = add_slide("OBJECT_WITH_CAPTION_TEXT")
set_title(s5, "Validation Framework")

body5L = set_body(s5, 2)
tf5L = clear_body(body5L)

p = tf5L.paragraphs[0]
r = p.add_run()
r.text = "Walk-Forward CV"
r.font.size = Pt(14)
r.font.bold = True
r.font.color.rgb = SMU_BLUE

add_bullet(tf5L, "756-day training window", font_size=12)
add_bullet(tf5L, "63-day test window", font_size=12)
add_bullet(tf5L, "5-day purge gap", font_size=12)

p2 = tf5L.add_paragraph()
p2.space_before = Pt(10)
r2 = p2.add_run()
r2.text = "Pre-Registered Hypotheses"
r2.font.size = Pt(14)
r2.font.bold = True
r2.font.color.rgb = SMU_BLUE

add_bullet(tf5L, "H1: Sentiment leads regime transitions", font_size=11)
add_bullet(tf5L, "H2: Divergence increases before transitions", font_size=11)
add_bullet(tf5L, "H3: Network connectedness rises pre-crisis", font_size=11)

remove_placeholder(s5, 1)
add_image_to_slide(s5, "fig_5_5_regime_vs_vix_covid_v2.png",
                   Inches(5.67), Inches(1.08), width=Inches(6.75))

# ══════════════════════════════════════════════════════════════
# SLIDE 6: Classification Performance (with confusion matrix)
# ══════════════════════════════════════════════════════════════
print("Creating Slide 6: Classification Performance")
s6 = add_slide("OBJECT_WITH_CAPTION_TEXT")
set_title(s6, "Classification Performance")

body6L = set_body(s6, 2)
tf6L = clear_body(body6L)

p = tf6L.paragraphs[0]
r = p.add_run()
r.text = "Key Metrics"
r.font.size = Pt(14)
r.font.bold = True
r.font.color.rgb = SMU_BLUE

metrics = [
    ("Accuracy", "92.84%"),
    ("F1 Weighted", "0.9236"),
    ("MCC", "0.8279"),
    ("Transition Acc.", "80.10%"),
]
for name, val in metrics:
    p_m = tf6L.add_paragraph()
    p_m.space_after = Pt(3)
    r_name = p_m.add_run()
    r_name.text = f"{name}:  "
    r_name.font.size = Pt(13)
    r_name.font.color.rgb = LIGHT_GRAY
    r_val = p_m.add_run()
    r_val.text = val
    r_val.font.size = Pt(16)
    r_val.font.bold = True
    r_val.font.color.rgb = GOLD

p_note = tf6L.add_paragraph()
p_note.space_before = Pt(10)
r_note = p_note.add_run()
r_note.text = "MCC is primary metric — accounts for class imbalance across 4 regime states"
r_note.font.size = Pt(10)
r_note.font.italic = True
r_note.font.color.rgb = LIGHT_GRAY

remove_placeholder(s6, 1)
add_image_to_slide(s6, "fig_5_5_transition_confusion_matrix.png",
                   Inches(5.67), Inches(1.08), height=Inches(5.0))

# ══════════════════════════════════════════════════════════════
# SLIDE 7: Hypothesis 1 — Sentiment Lead Time (with figure)
# ══════════════════════════════════════════════════════════════
print("Creating Slide 7: Hypothesis 1 — Sentiment Lead Time")
s7 = add_slide("OBJECT_WITH_CAPTION_TEXT")
set_title(s7, "H1: Sentiment Lead Time")

body7L = set_body(s7, 2)
tf7 = clear_body(body7L)

p = tf7.paragraphs[0]
r_tag = p.add_run()
r_tag.text = "H1 — NOT CONFIRMED"
r_tag.font.size = Pt(16)
r_tag.font.bold = True
r_tag.font.color.rgb = SMU_RED

add_bullet(tf7, "Does sentiment lead regime transitions by 1–5 days?", font_size=12, bold=True)

p_sub = tf7.add_paragraph()
p_sub.space_before = Pt(8)
r_sub = p_sub.add_run()
r_sub.text = "Results"
r_sub.font.size = Pt(14)
r_sub.font.bold = True
r_sub.font.color.rgb = SMU_BLUE

add_bullet(tf7, "Peak alignment at lag 0 — no 1-5 day lead detected globally", font_size=11)
add_bullet(tf7, "Only 1 event window confirmed (needed ≥2 for project-level support)", font_size=11, color=LIGHT_GRAY)
add_bullet(tf7, "COVID-19 showed localized lead, but not generalizable", font_size=11, color=LIGHT_GRAY)

p_key = tf7.add_paragraph()
p_key.space_before = Pt(8)
r_key = p_key.add_run()
r_key.text = "→  Standalone sentiment is not a reliable leading indicator"
r_key.font.size = Pt(11)
r_key.font.italic = True
r_key.font.color.rgb = GOLD

remove_placeholder(s7, 1)
add_image_to_slide(s7, "fig_5_5_h1_lead_time_summary_v2.png",
                   Inches(5.67), Inches(1.08), height=Inches(5.33))

# ══════════════════════════════════════════════════════════════
# SLIDE 8: Hypotheses 2 & 3 — Divergence & Network (with figures)
# ══════════════════════════════════════════════════════════════
print("Creating Slide 8: Hypotheses 2 & 3 — Divergence & Network")
s8_hyp = add_slide("OBJECT_WITH_CAPTION_TEXT")
set_title(s8_hyp, "H2 & H3: Divergence & Network Effects")

body8hL = set_body(s8_hyp, 2)
tf8h = clear_body(body8hL)

p = tf8h.paragraphs[0]
r2_tag = p.add_run()
r2_tag.text = "H2 — SUPPORTED  ✓"
r2_tag.font.size = Pt(14)
r2_tag.font.bold = True
r2_tag.font.color.rgb = RGBColor(0x2E, 0x8B, 0x57)

add_bullet(tf8h, "Cross-asset divergence increases before transitions", font_size=11)
add_bullet(tf8h, "Ratio = 1.19×  |  t = 18.73", font_size=10, color=LIGHT_GRAY)
add_bullet(tf8h, "p < 0.001  |  Cohen's d = 0.58", font_size=10, color=LIGHT_GRAY)

p3 = tf8h.add_paragraph()
p3.space_before = Pt(10)
r3_tag = p3.add_run()
r3_tag.text = "H3 — SUPPORTED  ✓"
r3_tag.font.size = Pt(14)
r3_tag.font.bold = True
r3_tag.font.color.rgb = RGBColor(0x2E, 0x8B, 0x57)

add_bullet(tf8h, "Network connectedness differs across regimes", font_size=11)
add_bullet(tf8h, "TCI: 0.5402 (stable) vs 0.4569 (trans.)", font_size=10, color=LIGHT_GRAY)
add_bullet(tf8h, "ANOVA F = 427.44  |  p < 0.001", font_size=10, color=LIGHT_GRAY)

p_key = tf8h.add_paragraph()
p_key.space_before = Pt(8)
r_key = p_key.add_run()
r_key.text = "→  Market disagreement (not level) signals regime shifts"
r_key.font.size = Pt(11)
r_key.font.italic = True
r_key.font.color.rgb = GOLD

remove_placeholder(s8_hyp, 1)
add_image_to_slide(s8_hyp, "fig_5_6_h2_divergence_distribution.png",
                   Inches(5.67), Inches(1.08), width=Inches(6.75))
add_image_to_slide(s8_hyp, "fig_5_6_h3_connectedness_comparison.png",
                   Inches(5.67), Inches(3.75), width=Inches(6.75))

# ══════════════════════════════════════════════════════════════
# SLIDE 9: Live Production Dashboard
# ══════════════════════════════════════════════════════════════
print("Creating Slide 9: Live Dashboard")
s9_dash = add_slide("OBJECT_WITH_CAPTION_TEXT")
set_title(s9_dash, "Live Production Dashboard")

body9dL = set_body(s9_dash, 2)
tf9dL = clear_body(body9dL)

p = tf9dL.paragraphs[0]
r = p.add_run()
r.text = "Backend"
r.font.size = Pt(14)
r.font.bold = True
r.font.color.rgb = SMU_BLUE

add_bullet(tf9dL, "FastAPI (Python)", font_size=12)
add_bullet(tf9dL, "PostgreSQL database", font_size=12)
add_bullet(tf9dL, "Railway hosting", font_size=12)
add_bullet(tf9dL, "Daily CRON refresh", font_size=12)

p2 = tf9dL.add_paragraph()
p2.space_before = Pt(10)
r2 = p2.add_run()
r2.text = "Frontend"
r2.font.size = Pt(14)
r2.font.bold = True
r2.font.color.rgb = SMU_BLUE

add_bullet(tf9dL, "Next.js + TypeScript", font_size=12)
add_bullet(tf9dL, "Recharts visualization", font_size=12)
add_bullet(tf9dL, "Vercel deployment", font_size=12)

p_url = tf9dL.add_paragraph()
p_url.space_before = Pt(10)
r_url = p_url.add_run()
r_url.text = "sentiment-regime-detector.vercel.app"
r_url.font.size = Pt(10)
r_url.font.italic = True
r_url.font.color.rgb = GOLD

remove_placeholder(s9_dash, 1)
add_image_to_slide(s9_dash, "dashboard_regime_overview.png",
                   Inches(5.67), Inches(1.08), width=Inches(6.75))

# ══════════════════════════════════════════════════════════════
# SLIDE 10: Discussion & Limitations
# ══════════════════════════════════════════════════════════════
print("Creating Slide 10: Discussion & Limitations")
s10_disc = add_slide("Content")
set_title(s10_disc, "Discussion & Limitations")
body10d = set_body(s10_disc, 1)
tf9 = clear_body(body10d)

p = tf9.paragraphs[0]
r = p.add_run()
r.text = "Key Findings"
r.font.size = Pt(16)
r.font.bold = True
r.font.color.rgb = SMU_BLUE

add_bullet(tf9, "Cross-asset sentiment divergence (H2) is the strongest predictive signal", font_size=13)
add_bullet(tf9, "Network effects (H3) confirm information cascades during regime transitions", font_size=13)
add_bullet(tf9, "Individual sentiment-to-regime lag (H1) is not reliable as standalone indicator", font_size=13)

p2 = tf9.add_paragraph()
p2.space_before = Pt(12)
r2 = p2.add_run()
r2.text = "Limitations"
r2.font.size = Pt(16)
r2.font.bold = True
r2.font.color.rgb = SMU_RED

add_bullet(tf9, "Computational cost — repeated CV tuning and rolling validation are expensive", font_size=13)
add_bullet(tf9, "Transformer inference latency — may limit near-real-time operation on large streams", font_size=13)
add_bullet(tf9, "Mixed-frequency alignment noise — irregular text vs. regular market bars", font_size=13)
add_bullet(tf9, "Data quality — bots, sarcasm, selection bias in public discourse", font_size=13)
add_bullet(tf9, "Correlational design — observed relationships may not be causal", font_size=13)

# ══════════════════════════════════════════════════════════════
# SLIDE 11: Conclusion & Future Work
# ══════════════════════════════════════════════════════════════
print("Creating Slide 11: Conclusion & Future Work")
s11 = add_slide("Content")
set_title(s11, "Conclusion & Future Work")
body11 = set_body(s11, 1)
tf10 = clear_body(body11)

p = tf10.paragraphs[0]
r = p.add_run()
r.text = "Contributions"
r.font.size = Pt(16)
r.font.bold = True
r.font.color.rgb = SMU_BLUE

add_bullet(tf10, "First unified framework combining NLP sentiment with volatility regime detection across asset classes", font_size=13)
add_bullet(tf10, "Novel cross-asset divergence metric as regime transition signal (H2)", font_size=13)
add_bullet(tf10, "Production-grade dashboard with daily automated updates", font_size=13)
add_bullet(tf10, "92.84% accuracy, 0.8279 MCC on 4-class regime classification", font_size=13)

p2 = tf10.add_paragraph()
p2.space_before = Pt(12)
r2 = p2.add_run()
r2.text = "Future Directions"
r2.font.size = Pt(16)
r2.font.bold = True
r2.font.color.rgb = SMU_BLUE

add_bullet(tf10, "Adaptive model weighting (IMCA-style recalibration)", font_size=13)
add_bullet(tf10, "Multimodal signals (earnings-call audio) and high-frequency decomposition", font_size=13)
add_bullet(tf10, "Expand to additional asset classes (bonds, REITs) and multilingual discourse", font_size=13)
add_bullet(tf10, "Causal identification to distinguish sentiment-to-price from reverse causality", font_size=13)

p3 = tf10.add_paragraph()
p3.space_before = Pt(14)
p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run()
r3.text = "Thank You — Questions?"
r3.font.size = Pt(20)
r3.font.bold = True
r3.font.color.rgb = GOLD

# ── Save ──
prs.save(OUTPUT)
print(f"\nPresentation saved to: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
